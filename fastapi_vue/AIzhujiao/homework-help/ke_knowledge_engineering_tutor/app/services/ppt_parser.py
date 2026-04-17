from io import BytesIO
from typing import Dict, List

from pptx import Presentation


def extract_ppt_texts(file_bytes: bytes) -> List[Dict]:
    prs = Presentation(BytesIO(file_bytes))
    slides = []
    for idx, slide in enumerate(prs.slides, start=1):
        texts = []
        title = ''
        for shape in slide.shapes:
            if not hasattr(shape, 'text'):
                continue
            text = (shape.text or '').strip()
            if not text:
                continue
            texts.append(text)
            if not title:
                title = text.splitlines()[0][:80]
        slides.append({
            'slide_no': idx,
            'title': title or f'第{idx}页',
            'text': '\n'.join(texts).strip(),
        })
    return slides
